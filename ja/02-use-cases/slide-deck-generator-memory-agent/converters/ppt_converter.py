"""
python-pptx を使用した HTML から PowerPoint への変換モジュール
"""

import os
import re
import uuid
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class HTMLToPowerPointConverter:
    """HTML スライドを PowerPoint プレゼンテーションに変換するクラス"""

    def __init__(self, output_dir: str):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

    def parse_html_slides(
        self, html_content: str
    ) -> Tuple[List[Dict[str, Any]], Dict[str, str]]:
        """HTML コンテンツを解析してスライドデータとスタイリングを抽出する"""
        soup = BeautifulSoup(html_content, "html.parser")
        slides = []
        styling = {}

        # Extract CSS variables for styling
        style_tag = soup.find("style")
        if style_tag:
            css_content = style_tag.get_text()
            styling = self.extract_css_variables(css_content)

        # Extract slides
        slide_divs = soup.find_all("div", class_="slide")

        for i, slide_div in enumerate(slide_divs):
            slide_content_div = slide_div.find("div", class_="slide-content")
            if not slide_content_div:
                continue

            slide_data = {"slide_number": i + 1}

            # Determine slide type
            if "title-slide" in slide_content_div.get("class", []):
                slide_data["type"] = "title"
                slide_data["title"] = self.extract_text(
                    slide_content_div.find(class_="main-title")
                )
                slide_data["subtitle"] = self.extract_text(
                    slide_content_div.find(class_="subtitle")
                )
                slide_data["author"] = self.extract_text(
                    slide_content_div.find(class_="author")
                )

            elif "content-slide" in slide_content_div.get("class", []):
                slide_data["type"] = "content"
                slide_data["title"] = self.extract_text(
                    slide_content_div.find(class_="slide-title")
                )

                # Extract bullet points
                bullet_list = slide_content_div.find(class_="bullet-list")
                if bullet_list:
                    slide_data["bullet_points"] = [
                        self.extract_text(li) for li in bullet_list.find_all("li")
                    ]

                # Extract slide text
                slide_text_div = slide_content_div.find(class_="slide-text")
                if slide_text_div:
                    slide_data["content"] = self.extract_text(slide_text_div)

            elif "section-slide" in slide_content_div.get("class", []):
                slide_data["type"] = "section"
                slide_data["title"] = self.extract_text(
                    slide_content_div.find(class_="section-title")
                )
                slide_data["subtitle"] = self.extract_text(
                    slide_content_div.find(class_="section-subtitle")
                )

            slides.append(slide_data)

        return slides, styling

    def extract_text(self, element) -> str:
        """HTML 要素からテキストを安全に抽出する"""
        return element.get_text(strip=True) if element else ""

    def extract_css_variables(self, css_content: str) -> Dict[str, str]:
        """スタイリング用の CSS カスタムプロパティ（変数）を抽出する"""
        variables = {}

        # Find :root block
        root_match = re.search(r":root\s*\{([^}]*)\}", css_content)
        if root_match:
            root_content = root_match.group(1)

            # Extract CSS variables
            var_pattern = r"--([^:]+):\s*([^;]+);"
            for match in re.finditer(var_pattern, root_content):
                var_name = match.group(1).strip()
                var_value = match.group(2).strip()
                variables[var_name] = var_value

        return variables

    def hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """16進数カラーを RGB タプルに変換する"""
        hex_color = hex_color.lstrip("#")
        if len(hex_color) != 6:
            return (0, 0, 0)  # Default to black

        try:
            return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))
        except ValueError:
            return (0, 0, 0)

    def apply_color_from_css(self, font, color_var: str, styling: Dict[str, str]):
        """CSS 変数からフォントに色を適用する"""
        if color_var in styling:
            hex_color = styling[color_var]
            r, g, b = self.hex_to_rgb(hex_color)
            font.color.rgb = RGBColor(r, g, b)

    def create_title_slide(
        self, prs: Presentation, slide_data: Dict[str, Any], styling: Dict[str, str]
    ):
        """PowerPoint でタイトルスライドを作成する"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        if slide_data.get("title"):
            title_shape = slide.shapes.title
            title_shape.text = slide_data["title"]

            # Apply styling
            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(36)
                    run.font.bold = True
                    self.apply_color_from_css(run.font, "primary-color", styling)

        # Set subtitle
        if slide_data.get("subtitle") and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = slide_data["subtitle"]

            for paragraph in subtitle_shape.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(24)
                    self.apply_color_from_css(run.font, "secondary-color", styling)

        # Add author at bottom if provided
        if slide_data.get("author"):
            left = Inches(1)
            top = Inches(7)
            width = Inches(8)
            height = Inches(0.5)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = slide_data["author"]

            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(14)
                    self.apply_color_from_css(run.font, "text-color", styling)

    def create_content_slide(
        self, prs: Presentation, slide_data: Dict[str, Any], styling: Dict[str, str]
    ):
        """PowerPoint でコンテンツスライドを作成する"""
        slide_layout = prs.slide_layouts[1]  # Content slide layout
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        if slide_data.get("title"):
            title_shape = slide.shapes.title
            title_shape.text = slide_data["title"]

            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(32)
                    run.font.bold = True
                    self.apply_color_from_css(run.font, "primary-color", styling)

        # Add content
        content_placeholder = None
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 2:  # Content placeholder
                content_placeholder = placeholder
                break

        if content_placeholder:
            text_frame = content_placeholder.text_frame
            text_frame.clear()

            # Add bullet points if available
            if slide_data.get("bullet_points"):
                for i, bullet_point in enumerate(slide_data["bullet_points"]):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()

                    p.text = bullet_point
                    p.level = 0

                    for run in p.runs:
                        run.font.size = Pt(20)
                        self.apply_color_from_css(run.font, "text-color", styling)

            # Add regular content if available
            elif slide_data.get("content"):
                p = text_frame.paragraphs[0]
                p.text = slide_data["content"]

                for run in p.runs:
                    run.font.size = Pt(18)
                    self.apply_color_from_css(run.font, "text-color", styling)

    def create_section_slide(
        self, prs: Presentation, slide_data: Dict[str, Any], styling: Dict[str, str]
    ):
        """PowerPoint でセクション区切りスライドを作成する"""
        slide_layout = prs.slide_layouts[
            2
        ]  # Section header layout (or use title layout)
        slide = prs.slides.add_slide(slide_layout)

        # Set section title
        if slide_data.get("title"):
            title_shape = slide.shapes.title
            title_shape.text = slide_data["title"]

            for paragraph in title_shape.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(40)
                    run.font.bold = True
                    self.apply_color_from_css(run.font, "primary-color", styling)

        # Set section subtitle
        if slide_data.get("subtitle") and len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = slide_data["subtitle"]

            for paragraph in subtitle_shape.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(24)
                    self.apply_color_from_css(run.font, "secondary-color", styling)

    def convert_html_to_pptx(
        self, html_file_path: str, output_filename: Optional[str] = None
    ) -> str:
        """HTML プレゼンテーションを PowerPoint に変換する"""

        # Read HTML file
        with open(html_file_path, "r", encoding="utf-8") as f:
            html_content = f.read()

        # Parse slides and styling
        slides_data, styling = self.parse_html_slides(html_content)

        if not slides_data:
            raise ValueError("No slides found in HTML file")

        # Create PowerPoint presentation
        prs = Presentation()

        # Apply theme colors if available
        if styling:
            self.apply_theme_colors(prs, styling)

        # Create slides
        for slide_data in slides_data:
            slide_type = slide_data.get("type", "content")

            if slide_type == "title":
                self.create_title_slide(prs, slide_data, styling)
            elif slide_type == "content":
                self.create_content_slide(prs, slide_data, styling)
            elif slide_type == "section":
                self.create_section_slide(prs, slide_data, styling)

        # Save PowerPoint file
        if output_filename is None:
            base_name = os.path.splitext(os.path.basename(html_file_path))[0]
            output_filename = f"{base_name}_{uuid.uuid4().hex[:8]}.pptx"

        output_path = os.path.join(self.output_dir, output_filename)
        prs.save(output_path)

        return output_path

    def apply_theme_colors(self, prs: Presentation, styling: Dict[str, str]):
        """プレゼンテーションにテーマカラーを適用する（python-pptx でのサポートは限定的）"""
        # 注: python-pptx はテーマ変更のサポートが限定的
        # これはデフォルトカラーを設定する基本的な実装
        pass

    def convert_slides_data_to_pptx(
        self,
        slides_data: List[Dict[str, Any]],
        title: str = "Presentation",
        styling: Optional[Dict[str, str]] = None,
    ) -> str:
        """HTML ファイルを経由せずにスライドデータを直接 PowerPoint に変換する"""

        if not slides_data:
            raise ValueError("No slide data provided")

        if styling is None:
            styling = {}

        # Create PowerPoint presentation
        prs = Presentation()

        # Create slides
        for slide_data in slides_data:
            slide_type = slide_data.get("type", "content")

            if slide_type == "title":
                self.create_title_slide(prs, slide_data, styling)
            elif slide_type == "content":
                self.create_content_slide(prs, slide_data, styling)
            elif slide_type == "section":
                self.create_section_slide(prs, slide_data, styling)

        # Save PowerPoint file
        output_filename = f"{title.replace(' ', '_')}_{uuid.uuid4().hex[:8]}.pptx"
        output_path = os.path.join(self.output_dir, output_filename)
        prs.save(output_path)

        return output_path


# 使用例とテスト
def demo_html_to_ppt_converter():
    """HTML から PowerPoint への変換をデモンストレーションする"""
    converter = HTMLToPowerPointConverter("output")

    # Sample slide data for testing
    sample_slides = [
        {
            "type": "title",
            "title": "Introduction to AI",
            "subtitle": "A Comprehensive Overview",
            "author": "Generated on " + datetime.now().strftime("%B %d, %Y"),
        },
        {
            "type": "content",
            "title": "What is Artificial Intelligence?",
            "bullet_points": [
                "Computer systems that can perform tasks requiring human intelligence",
                "Includes machine learning, natural language processing, and computer vision",
                "Applications in healthcare, finance, transportation, and more",
            ],
        },
        {
            "type": "section",
            "title": "Types of AI",
            "subtitle": "Understanding Different Approaches",
        },
        {
            "type": "content",
            "title": "Machine Learning Categories",
            "bullet_points": [
                "Supervised Learning: Learning from labeled examples",
                "Unsupervised Learning: Finding patterns in unlabeled data",
                "Reinforcement Learning: Learning through rewards and penalties",
            ],
        },
    ]

    sample_styling = {
        "primary-color": "#2563eb",
        "secondary-color": "#3b82f6",
        "text-color": "#1e293b",
        "accent-color": "#0ea5e9",
    }

    try:
        pptx_path = converter.convert_slides_data_to_pptx(
            sample_slides, "AI Introduction Demo", sample_styling
        )
        print(f"✅ PowerPoint プレゼンテーションを作成しました: {pptx_path}")
        return pptx_path
    except Exception as e:
        print(f"❌ PowerPoint の作成でエラーが発生しました: {e}")
        return None


if __name__ == "__main__":
    demo_html_to_ppt_converter()
